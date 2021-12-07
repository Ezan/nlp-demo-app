from functools import reduce
import os
import sys
from pptx.dml.color import ColorFormat, RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt
from numpy.core.fromnumeric import shape
from pptx import Presentation
from pptx.oxml.xmlchemy import OxmlElement
from pptx.enum.shapes import PP_PLACEHOLDER, MSO_SHAPE_TYPE
# import re
import win32com.client as win32
import pandas as pd



# 9 slide layouts
if __name__ == "__main__":
    project_root_path = os.path.split(os.path.dirname(os.path.realpath(__file__)))[0]
    sys.path.insert(0, project_root_path)

class PPTReader:

    def __init__(self, filename):
        # self.text_runs = []
        if (filename.split('.')[-1] == 'ppt'):
            Application = win32.Dispatch("PowerPoint.Application")
            Application.Visible = True
            prst = Application.Presentations.Open(os.path.join(os.curdir,"input/{}".format(filename)))
            prst.Saveas(os.path.join(os.curdir,"input/{}".format("temp.pptx")))
            filename = "temp.pptx"
            prst.Close()
            Application.Quit()
        self.pathToPresentation = os.path.join(os.curdir,"input/{}".format(filename))
        self.presentation = Presentation(self.pathToPresentation)
        self._textStore = {}
        self.largestFont = []
        self._tobeSkippedPlaceholders = [PP_PLACEHOLDER.TITLE,PP_PLACEHOLDER.PICTURE, PP_PLACEHOLDER.BITMAP, PP_PLACEHOLDER.SUBTITLE, PP_PLACEHOLDER.CENTER_TITLE, PP_PLACEHOLDER.CHART,
            PP_PLACEHOLDER.DATE, PP_PLACEHOLDER.FOOTER, PP_PLACEHOLDER.HEADER,PP_PLACEHOLDER.MEDIA_CLIP, PP_PLACEHOLDER.ORG_CHART, PP_PLACEHOLDER.TABLE, PP_PLACEHOLDER.VERTICAL_TITLE]

    def highlightpptbyresult(self, results):
        for result in results:
            slide = self.presentation.slides.get(result['slide_id'])
            self.__recursiveShapeHighlight(slide.shapes, result)
        return self.presentation
            


    def __recursiveShapeHighlight(self, shapes, result):
        for shape in shapes:
            if shape.is_placeholder:
                if shape.placeholder_format.type in self._tobeSkippedPlaceholders:
                    continue
            # if  shape.is_placeholder and (not shape.has_text_frame or shape.has_chart):
            #     continue
            else:
                if (shape.shape_type not in [MSO_SHAPE_TYPE.GROUP, MSO_SHAPE_TYPE.TEXT_BOX]):
                    continue    
                elif (shape.shape_type is MSO_SHAPE_TYPE.GROUP):
                    self.__recursiveShapeHighlight(shape.shapes, result)
                if (not shape.has_text_frame):
                    continue
                elif len(shape.text_frame.text) < 15:
                    continue
            frame = shape.text_frame.text
            if len(frame.split()) < 5:
                continue    
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if result['result'].split(" ")[0].lower() in run.text.lower().split(" ") or result['result'].split(" ")[-1].lower() in run.text.lower().split(" ") :
                    # if run.text.lower().find(result['result'].lower()) != -1 and result['result'].lower().find(run.text.lower()) != -1:
                        run.font.color.rgb = RGBColor(255, 50, 50)


    def __recursiveShapeParser(self, shapes, slide_idx,slide_id,slide_obj):
        for shape_idx, shape in enumerate(shapes):
            if shape.is_placeholder:
                if shape.placeholder_format.type in self._tobeSkippedPlaceholders:
                    continue
            # if  shape.is_placeholder and (not shape.has_text_frame or shape.has_chart):
            #     continue
            else:
                if (shape.shape_type not in [MSO_SHAPE_TYPE.GROUP, MSO_SHAPE_TYPE.TEXT_BOX]):
                    continue
                elif (shape.shape_type is MSO_SHAPE_TYPE.GROUP):
                    self.__recursiveShapeParser(shape.shapes, slide_idx, slide_id, slide_obj)
                if (not shape.has_text_frame):
                    continue
                elif len(shape.text_frame.text) < 15:
                    continue
            frame = shape.text_frame.text
            if len(frame.split()) < 5:
                continue            
            self.largestFont += (list(map(lambda x: { 'slide_idx':slide_idx, 'slide_id':slide_id,
            'shape_idx':shape_idx, 'shape_id':shape.shape_id, 'fontInfo': x }, reduce(self.__fixer, shape.text_frame.paragraphs, []))))
            slide_obj[slide_idx]["shapes"][shape_idx] = {"shape_idx": shape_idx, "shape_id": shape.shape_id, "text_frame": shape.text_frame }
        return slide_obj

    def __fixer(self, maxTuple, item):
        lot = list(map(lambda p: (p.font.size, p.font.name), item.runs))
        from operator import itemgetter
        if len(lot) > 0:
            maxTuple.append(max(lot, key=itemgetter(0)))
        return maxTuple

    def readText(self):
        print("Inside ppt File")
        for slide_idx, slide in enumerate(self.presentation.slides):
            slide_obj = { slide_idx : { "slide_id": slide.slide_id, "shapes": {} } }
            slide_obj = self.__recursiveShapeParser(shapes=slide.shapes, slide_idx=slide_idx, slide_id=slide.slide_id, slide_obj=slide_obj)
            # for shape_idx, shape in enumerate(slide.shapes):
            #     if not shape.has_text_frame:
            #         continue
            #     for paragraph in shape.text_frame.paragraphs:
            #         for run in paragraph.runs:
            #             self.text_runs.append({'shape_id': shape_idx, 'text': run.text})

            self._textStore = { **self._textStore, **slide_obj}

    def fetchPresentation(self):
        return self.presentation
        
    def fetchTextStore(self):
        return self._textStore

    def fontWiseFilter(self):
        frame = pd.concat([pd.DataFrame(self.largestFont, columns=['slide_idx','slide_id', 'shape_idx','shape_id']),pd.DataFrame(list(map(lambda x: x['fontInfo'], self.largestFont)), columns=['font_size', 'font_name'])], axis=1)
        print(frame.head())
        df_tmp = frame.groupby(['slide_idx','slide_id']).agg({'shape_idx': 'first', 'shape_id': 'first', 'font_size':'max'}).reset_index()
        # df_tmp['occurance'] = frame.groupby(['slide_idx', 'slide_id'])['font_size'].transform('count')
        df_tmp = df_tmp.reset_index()
        df_tmp.drop_duplicates(inplace=True)
        newFrame = frame.groupby(['slide_idx', 'slide_id'])['font_size'].value_counts()
        filteredDict = [(s , k , v) for (s,k,v), val in newFrame.to_dict().items() if val <= 2]
        df_tmp['tmp'] = list(zip(df_tmp['slide_idx'],df_tmp['slide_id'],df_tmp['font_size']))
        df_tmp = df_tmp[df_tmp['tmp'].isin(filteredDict)]
        for (s,ss) in list(df_tmp[['slide_idx','shape_idx']].itertuples(index=False, name=None)):
            self._textStore[s]['shapes'].pop(ss)

if __name__ == "__main__":
    # print("Invoked readFromPPT")
    name = "sample1.pptx"
    reader = PPTReader(name)
    reader.readText()
    reader.displayText()
    store = reader.fetchTextStore()
    reader.fontWiseFilter()
    # from tika import parser
    # parsed = parser.from_file(os.path.join(os.curdir,"pptFiles/{}".format("sample1.pptx")))
    # print(parsed)
    # print(store)